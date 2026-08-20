"""
Document generators — produce PDF, DOCX, XLSX, PPTX, Markdown, and HTML→PDF
files from a structured content schema, and register them as attachments
via the file_storage backend.

Each generator returns an Attachment dataclass. The tool executor wraps
these into the per-call pending_attachments list; agent_runner drains that
list when persisting the assistant message so the frontend receives an
`attachment` WS event and the DocumentSplit pane can render the file.

Heavy deps (reportlab, openpyxl, python-docx, python-pptx, weasyprint)
are imported lazily at call-site so the platform image (which doesn't
need weasyprint) can skip installing Pango/Cairo.
"""

from __future__ import annotations

import asyncio
import io
import logging
import os
import re
import unicodedata
import uuid
from dataclasses import dataclass, asdict
from datetime import datetime, timezone
from typing import Any, Dict, Iterable, List, Optional

from app.services.file_storage import get_storage_backend


logger = logging.getLogger(__name__)

# Strong refs to in-flight preview warms — see _prewarm_preview.
_PREWARM_TASKS: set[asyncio.Task] = set()


class EmptyDocumentError(ValueError):
    """A generator was asked to build a document with nothing in it.

    Raised BEFORE anything is persisted, so an empty call leaves no file on
    disk and nothing to attach. Founder incident 2026-08-18: the model
    opened a research turn with `generate_pdf({"filename": "x", "content":
    []})` — reportlab happily built a 952-byte PDF containing only the page
    footer, `_persist` stored it, `_register_attachment` attached it, and
    the user received an empty `x.pdf` they never asked for. The message
    is written for the model that reads the tool result: it says what was
    NOT done and what to do instead, so a stub call ends in a correction,
    not a fabricated "here's your file".
    """

    def __init__(self, tool: str, what: str):
        super().__init__(
            f"{tool} was called with no {what} — nothing was created or "
            f"attached. If the user asked for a document, call again with "
            f"the real content; if they did not, answer in chat instead."
        )
        self.tool = tool


MIME_PDF = "application/pdf"
MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_MD = "text/markdown"
MIME_HTML = "text/html"


@dataclass
class Attachment:
    id: str
    filename: str
    mime_type: str
    size_bytes: int
    storage_path: str
    created_at: str
    # Intrinsic pixels for image attachments; None for everything else and for
    # an image whose header would not decode. Clients lay the card out at this
    # shape on FIRST paint — without it they have to guess an aspect ratio and
    # the picture jumps when it decodes. These ride the Message.attachments JSON
    # column, so carrying them needed no migration; old rows simply answer None
    # and the client falls back to measuring the file itself.
    width: Optional[int] = None
    height: Optional[int] = None

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


def _safe_filename(name: str, default_ext: str) -> str:
    """Strip directory traversal and ensure the requested extension."""
    # Normalize Windows-style separators BEFORE basename — on POSIX,
    # os.path.basename("..\\..\\etc\\x") returns the whole string, which
    # would otherwise survive into the storage key as a literal filename.
    name = (name or "").replace("\\", "/")
    name = os.path.basename(name).strip() or f"document{default_ext}"
    if not name.lower().endswith(default_ext):
        name = f"{os.path.splitext(name)[0] or 'document'}{default_ext}"
    return name


def _prewarm_preview(storage_key: str, mime_type: str) -> None:
    """Start the DOCX/PPTX -> PDF conversion now, not on first view.

    The preview pane converts through LibreOffice, and that conversion was
    measured at 10.3 s cold (3.3 s warm, 0.2 s cached) on a real user file
    on 2026-08-13. Every generated document is cold exactly once — on the
    view that immediately follows generation, which is the view the user
    always takes. So the one time it matters is the one time it is slow.

    Doing it here costs nothing the user waits for: the agent has just
    written the bytes and is still composing its reply, so the conversion
    overlaps with the rest of the turn and the pane opens from cache.

    Fire-and-forget by design. A failure here must never fail the tool
    call that produced the document — the preview route still converts on
    demand, and the user gets a file either way.
    """
    if mime_type not in (MIME_DOCX, MIME_PPTX):
        return
    try:
        loop = asyncio.get_running_loop()
    except RuntimeError:
        return

    async def _run() -> None:
        try:
            from app.services.doc_preview import render_to_pdf
            await render_to_pdf(storage_key)
        except Exception as e:  # never surface — this is pure latency work
            logger.debug("preview prewarm skipped for %s: %s", storage_key, e)

    task = loop.create_task(_run())
    # Hold a reference; a bare create_task can be garbage-collected
    # mid-flight, which would silently drop the warm we just scheduled.
    _PREWARM_TASKS.add(task)
    task.add_done_callback(_PREWARM_TASKS.discard)


# ── Meaningful filenames ─────────────────────────────────────────
#
# The system prompt asks for descriptive filenames ("march-expenses.xlsx,
# not file.xlsx") but nothing enforced it: `generate_pdf(filename="x")`
# shipped as `x.pdf` (2026-08-18). A placeholder stem is replaced by the
# best available hint — the document title, its first heading / sheet /
# slide title, the first line of prose — and, failing all of those, a
# dated `document-YYYY-MM-DD` name. A descriptive stem the model chose is
# always kept verbatim.

# Stems that carry no information. Matched case-insensitively against the
# basename without extension. A single character or an all-digit stem is a
# placeholder too (see _is_placeholder_stem). Two-letter stems are NOT
# blanket-rejected — "cv", "q3", "ai" are real names — only the ones
# listed here.
_PLACEHOLDER_STEMS = frozenset({
    "x", "xx", "xxx", "yy", "zz", "aa", "ab", "abc", "asdf", "test", "testing",
    "tmp", "temp", "file", "files", "doc", "docs", "document", "documents",
    "untitled", "new", "output", "out", "result", "results", "generated",
    "attachment", "download", "foo", "bar", "baz", "placeholder", "sample",
    "example", "none", "null", "undefined", "unnamed", "noname", "default",
    "pdf", "docx", "xlsx", "pptx", "md", "markdown", "txt", "html",
    "sheet", "sheet1", "workbook", "book", "book1", "slides", "deck",
    "presentation", "spreadsheet", "export", "data",
})

_SLUG_STRIP_RE = re.compile(r"[^a-z0-9]+")


def _is_placeholder_stem(stem: str) -> bool:
    s = (stem or "").strip().lower()
    if len(s) <= 1:
        return True
    if s.isdigit():
        return True
    # "file1", "test2", "document (3)", "untitled-1" — a placeholder with a
    # counter or decoration is still a placeholder.
    core = _SLUG_STRIP_RE.sub("", s).rstrip("0123456789")
    return s in _PLACEHOLDER_STEMS or core in _PLACEHOLDER_STEMS


def _slugify(text: Any, max_len: int = 60) -> str:
    """'March 2026 — Expenses & Totals' → 'march-2026-expenses-totals'.
    ASCII-folds accents; anything non-alphanumeric becomes a hyphen; the
    result is trimmed at a hyphen boundary so a long title never ends in a
    truncated word. Returns '' when nothing survives (e.g. all-CJK input —
    the caller falls back to the dated name rather than transliterating)."""
    if text is None:
        return ""
    s = str(text)
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode("ascii")
    s = _SLUG_STRIP_RE.sub("-", s.lower()).strip("-")
    if len(s) > max_len:
        s = s[:max_len].rsplit("-", 1)[0] or s[:max_len]
    return s.strip("-")


def _first_words(text: Any, n: int = 6) -> str:
    """First ``n`` words of a block of prose, for a filename hint."""
    words = str(text or "").split()
    return " ".join(words[:n])


def _derive_filename(
    requested: Optional[str],
    default_ext: str,
    *,
    title: Optional[str] = None,
    hints: Iterable[Any] = (),
) -> str:
    """The filename a generator should persist under.

    ``requested`` (the model's `filename` arg) wins whenever its stem is
    descriptive. Otherwise the first non-empty slug from ``title`` then
    ``hints`` (in order) is used, and if none yields a slug the file is
    named ``document-YYYY-MM-DD``. Always returns a safe basename with
    ``default_ext``.
    """
    safe = _safe_filename(requested or "", default_ext)
    stem = os.path.splitext(safe)[0]
    if not _is_placeholder_stem(stem):
        return safe
    for candidate in (title, *hints):
        slug = _slugify(candidate)
        if slug and not _is_placeholder_stem(slug):
            return f"{slug}{default_ext}"
    return f"document-{datetime.now(timezone.utc).strftime('%Y-%m-%d')}{default_ext}"


# ── Empty-content detection ──────────────────────────────────────


def _block_has_content(block: Any) -> bool:
    """True if a structured block would put something visible on the page.
    A `page_break` never counts; a table needs at least one header or row;
    a list needs an item; an image needs a path that exists."""
    if not isinstance(block, dict):
        return bool(str(block or "").strip())
    btype = block.get("type", "paragraph")
    if btype in ("heading", "paragraph"):
        return bool(str(block.get("text", "") or "").strip())
    if btype == "table":
        headers = [h for h in (block.get("headers") or []) if str(h or "").strip()]
        rows = [r for r in (block.get("rows") or []) if r]
        return bool(headers or rows)
    if btype in ("bullet_list", "numbered_list"):
        return any(str(i or "").strip() for i in (block.get("items") or []))
    if btype == "image":
        path = block.get("path", "")
        return bool(path) and os.path.isfile(path)
    if btype == "page_break":
        return False
    # Unknown block type — count it if it carries any text.
    return bool(str(block.get("text", "") or "").strip())


def _blocks_have_content(content: Any) -> bool:
    return any(_block_has_content(b) for b in _iter_blocks(content))


def _first_block_text(content: Any, *, kinds: tuple = ("heading",)) -> str:
    """Text of the first block of one of ``kinds`` (filename hint)."""
    for block in _iter_blocks(content):
        if isinstance(block, dict) and block.get("type", "paragraph") in kinds:
            text = str(block.get("text", "") or "").strip()
            if text:
                return text
    return ""


def _first_prose(content: Any) -> str:
    return _first_words(_first_block_text(content, kinds=("paragraph",)))


_HTML_TAG_RE = re.compile(r"<[^>]+>")
_HTML_TITLE_RE = re.compile(r"<title[^>]*>(.*?)</title>", re.IGNORECASE | re.DOTALL)
_HTML_H1_RE = re.compile(r"<h1[^>]*>(.*?)</h1>", re.IGNORECASE | re.DOTALL)


def _html_text(html: str) -> str:
    return _HTML_TAG_RE.sub(" ", html or "").strip()


def _image_dimensions(data: bytes, mime_type: str) -> tuple[Optional[int], Optional[int]]:
    """Intrinsic (width, height) for image bytes, or (None, None).

    Reads the HEADER only — Pillow is lazy, so `.size` never decodes the pixels.
    Deliberately lives in `_persist` rather than in the image tools: EVERY image
    that reaches a chat card goes through here (generate_image, edit_image,
    send_photo, convert_document, anything added later), so the create path and
    the edit path cannot end up stamping different things. Best-effort by
    design — a picture is worth delivering even when its header will not parse.
    """
    if not (mime_type or "").startswith("image/"):
        return None, None
    try:
        from PIL import Image as _PILImage  # already a dependency (HEIC decode)
        with _PILImage.open(io.BytesIO(data)) as im:
            w, h = im.size
        return (int(w), int(h)) if w > 0 and h > 0 else (None, None)
    except Exception:
        logger.debug("attachment: could not read image dimensions", exc_info=True)
        return None, None


async def _persist(data: bytes, filename: str, mime_type: str, user_scope: str) -> Attachment:
    """Write bytes to storage under {user_scope}/{uuid}_{filename} and return an Attachment."""
    att_id = uuid.uuid4().hex
    key = f"{user_scope}/{att_id}_{filename}" if user_scope else f"{att_id}_{filename}"
    backend = get_storage_backend()
    await backend.put(key, data)
    _prewarm_preview(key, mime_type)
    width, height = _image_dimensions(data, mime_type)
    return Attachment(
        id=att_id,
        filename=filename,
        mime_type=mime_type,
        size_bytes=len(data),
        storage_path=key,
        created_at=datetime.now(timezone.utc).isoformat(),
        width=width,
        height=height,
    )


# ── Structured content iterators shared by PDF/DOCX/PPTX ──────────

def _iter_blocks(content: Any) -> List[Dict[str, Any]]:
    """Normalize content to a list of block dicts.

    Accepts either a list of blocks, or a string (treated as a single
    paragraph). Blocks have a `type` key: heading|paragraph|table|image|page_break.
    """
    if isinstance(content, str):
        return [{"type": "paragraph", "text": content}]
    if isinstance(content, list):
        return content
    return []


# ── PDF via reportlab ─────────────────────────────────────────────

async def gen_pdf(
    content: Any,
    filename: str,
    *,
    user_scope: str,
    title: Optional[str] = None,
    cover_page: bool = False,
) -> Attachment:
    from reportlab.lib.pagesizes import LETTER  # type: ignore
    from reportlab.lib.styles import getSampleStyleSheet  # type: ignore
    from reportlab.platypus import (  # type: ignore
        SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image,
    )
    from reportlab.lib import colors  # type: ignore
    from reportlab.lib.units import inch  # type: ignore

    # Refuse before touching reportlab: an empty story still builds a
    # valid 952-byte PDF (page footer only), which is exactly what shipped
    # as x.pdf. A cover page with a title but no body is still empty.
    if not _blocks_have_content(content):
        raise EmptyDocumentError("generate_pdf", "content blocks")
    filename = _derive_filename(
        filename, ".pdf",
        title=title,
        hints=(_first_block_text(content), _first_prose(content)),
    )
    buf = io.BytesIO()

    def _on_page(canvas, doc):
        canvas.saveState()
        canvas.setFont("Helvetica", 9)
        canvas.setFillGray(0.4)
        canvas.drawRightString(LETTER[0] - 0.5 * inch, 0.4 * inch, f"Page {doc.page}")
        canvas.restoreState()

    doc = SimpleDocTemplate(buf, pagesize=LETTER, title=title or filename)
    styles = getSampleStyleSheet()
    story: list = []

    if cover_page and title:
        story.extend([
            Spacer(1, 3 * inch),
            Paragraph(f"<para align='center'><font size='24'><b>{title}</b></font></para>", styles["Title"]),
            Spacer(1, 0.3 * inch),
            Paragraph(
                f"<para align='center'>{datetime.now().strftime('%B %d, %Y')}</para>",
                styles["Normal"],
            ),
            PageBreak(),
        ])

    for block in _iter_blocks(content):
        btype = block.get("type", "paragraph")
        if btype == "heading":
            level = int(block.get("level", 1))
            style_name = {1: "Heading1", 2: "Heading2", 3: "Heading3", 4: "Heading4"}.get(level, "Heading2")
            story.append(Paragraph(block.get("text", ""), styles[style_name]))
        elif btype == "paragraph":
            story.append(Paragraph(block.get("text", ""), styles["BodyText"]))
            story.append(Spacer(1, 0.1 * inch))
        elif btype == "table":
            rows = [block.get("headers", [])] + list(block.get("rows", []))
            if rows and rows[0]:
                t = Table(rows, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]))
                story.append(t)
                story.append(Spacer(1, 0.15 * inch))
        elif btype == "image":
            path = block.get("path", "")
            if path and os.path.isfile(path):
                try:
                    story.append(Image(path, width=5 * inch, height=3 * inch, kind="proportional"))
                    if block.get("caption"):
                        story.append(Paragraph(
                            f"<para align='center'><i>{block['caption']}</i></para>",
                            styles["Italic"],
                        ))
                    story.append(Spacer(1, 0.15 * inch))
                except Exception:
                    pass
        elif btype == "page_break":
            story.append(PageBreak())

    doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    return await _persist(buf.getvalue(), filename, MIME_PDF, user_scope)


# ── DOCX via python-docx ──────────────────────────────────────────

async def gen_docx(
    content: Any,
    filename: str,
    *,
    user_scope: str,
    title: Optional[str] = None,
) -> Attachment:
    from docx import Document  # type: ignore
    from docx.shared import Inches  # type: ignore

    if not _blocks_have_content(content):
        raise EmptyDocumentError("generate_docx", "content blocks")
    filename = _derive_filename(
        filename, ".docx",
        title=title,
        hints=(_first_block_text(content), _first_prose(content)),
    )
    doc = Document()

    if title:
        doc.add_heading(title, level=0)

    for block in _iter_blocks(content):
        btype = block.get("type", "paragraph")
        if btype == "heading":
            level = max(1, min(4, int(block.get("level", 1))))
            doc.add_heading(block.get("text", ""), level=level)
        elif btype == "paragraph":
            doc.add_paragraph(block.get("text", ""))
        elif btype == "bullet_list":
            for item in block.get("items", []):
                doc.add_paragraph(str(item), style="List Bullet")
        elif btype == "numbered_list":
            for item in block.get("items", []):
                doc.add_paragraph(str(item), style="List Number")
        elif btype == "table":
            headers = block.get("headers", [])
            rows = block.get("rows", [])
            if headers:
                tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                tbl.style = "Light Grid Accent 1"
                for i, h in enumerate(headers):
                    tbl.rows[0].cells[i].text = str(h)
                for r, row in enumerate(rows, start=1):
                    for c, cell in enumerate(row[: len(headers)]):
                        tbl.rows[r].cells[c].text = "" if cell is None else str(cell)
        elif btype == "image":
            path = block.get("path", "")
            if path and os.path.isfile(path):
                try:
                    doc.add_picture(path, width=Inches(5))
                except Exception:
                    pass
        elif btype == "page_break":
            doc.add_page_break()

    buf = io.BytesIO()
    doc.save(buf)
    return await _persist(buf.getvalue(), filename, MIME_DOCX, user_scope)


# ── XLSX via openpyxl ─────────────────────────────────────────────

async def gen_xlsx(
    sheets: List[Dict[str, Any]],
    filename: str,
    *,
    user_scope: str,
) -> Attachment:
    from openpyxl import Workbook  # type: ignore
    from openpyxl.styles import Font, PatternFill  # type: ignore
    from openpyxl.utils import get_column_letter  # type: ignore

    sheets = [s for s in (sheets or []) if isinstance(s, dict)]
    if not any(
        [h for h in (s.get("headers") or []) if str(h or "").strip()]
        or [r for r in (s.get("rows") or []) if r]
        for s in sheets
    ):
        raise EmptyDocumentError("generate_xlsx", "sheet data (no headers, no rows)")
    filename = _derive_filename(
        filename, ".xlsx",
        hints=(
            next((s.get("name") for s in sheets if s.get("name")), None),
            " ".join(str(h) for h in (sheets[0].get("headers") or [])[:4]),
        ),
    )
    wb = Workbook()
    wb.remove(wb.active)  # Drop the default sheet

    for sheet_def in sheets:
        name = str(sheet_def.get("name", "Sheet"))[:31] or "Sheet"
        ws = wb.create_sheet(title=name)
        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])

        if headers:
            ws.append(headers)
            bold = Font(bold=True)
            fill = PatternFill("solid", fgColor="F0F0F0")
            for col_idx in range(1, len(headers) + 1):
                c = ws.cell(row=1, column=col_idx)
                c.font = bold
                c.fill = fill

        for row in rows:
            ws.append(list(row))

        # Auto-size columns (rough estimate — openpyxl doesn't support true auto-fit).
        for col_idx in range(1, max(1, len(headers)) + 1):
            letter = get_column_letter(col_idx)
            max_len = len(str(headers[col_idx - 1])) if col_idx - 1 < len(headers) else 8
            for row in rows:
                if col_idx - 1 < len(row):
                    v = row[col_idx - 1]
                    if v is not None:
                        max_len = max(max_len, len(str(v)))
            ws.column_dimensions[letter].width = min(60, max_len + 2)

    buf = io.BytesIO()
    wb.save(buf)
    return await _persist(buf.getvalue(), filename, MIME_XLSX, user_scope)


# ── PPTX via python-pptx ──────────────────────────────────────────

async def gen_pptx(
    slides: List[Dict[str, Any]],
    filename: str,
    *,
    user_scope: str,
) -> Attachment:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    slides = [s for s in (slides or []) if isinstance(s, dict)]

    def _slide_has_content(s: Dict[str, Any]) -> bool:
        if str(s.get("title", "") or "").strip() or str(s.get("subtitle", "") or "").strip():
            return True
        if any(str(b or "").strip() for b in (s.get("bullets") or [])):
            return True
        path = s.get("path", "")
        return bool(path) and os.path.isfile(path)

    if not any(_slide_has_content(s) for s in slides):
        raise EmptyDocumentError("generate_pptx", "slides")
    filename = _derive_filename(
        filename, ".pptx",
        hints=(next((s.get("title") for s in slides if str(s.get("title", "") or "").strip()), None),),
    )
    prs = Presentation()

    for slide_def in slides:
        stype = slide_def.get("type", "content")
        if stype == "title":
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_def.get("subtitle", "")
        elif stype == "section":
            layout = prs.slide_layouts[2] if len(prs.slide_layouts) > 2 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
        elif stype == "image":
            layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
            path = slide_def.get("path", "")
            if path and os.path.isfile(path):
                try:
                    slide.shapes.add_picture(path, Inches(1), Inches(1.5), width=Inches(8))
                except Exception:
                    pass
        else:  # "content" default
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_def.get("title", "")
            body = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
            bullets = slide_def.get("bullets", [])
            if body and bullets:
                body.text = str(bullets[0])
                for b in bullets[1:]:
                    p = body.add_paragraph()
                    p.text = str(b)
                    p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    return await _persist(buf.getvalue(), filename, MIME_PPTX, user_scope)


# ── Markdown (plain write) ────────────────────────────────────────

async def gen_markdown(content: str, filename: str, *, user_scope: str) -> Attachment:
    text = content if isinstance(content, str) else ("" if content is None else str(content))
    if not text.strip():
        raise EmptyDocumentError("generate_markdown", "content")
    # Hint: the first ATX heading, else the first non-blank line.
    first_heading = next(
        (ln.lstrip("#").strip() for ln in text.splitlines()
         if ln.lstrip().startswith("#") and ln.lstrip("#").strip()),
        "",
    )
    first_line = next((ln.strip() for ln in text.splitlines() if ln.strip()), "")
    filename = _derive_filename(
        filename, ".md", hints=(first_heading, _first_words(first_line)),
    )
    data = text.encode("utf-8")
    return await _persist(data, filename, MIME_MD, user_scope)


# ── HTML → PDF via weasyprint (agent-side only) ───────────────────

async def gen_html_to_pdf(html: str, filename: str, *, user_scope: str) -> Attachment:
    # Empty check first — a missing weasyprint must not mask a stub call,
    # and a stub call must not be reported as a missing renderer.
    html = html if isinstance(html, str) else ("" if html is None else str(html))
    if not _html_text(html):
        raise EmptyDocumentError("generate_html_to_pdf", "html content")
    try:
        from weasyprint import HTML  # type: ignore
    except ImportError as e:
        raise RuntimeError(
            "weasyprint is not installed in this runtime. "
            "HTML→PDF is agent-side only (requires Pango/Cairo). "
            "Use generate_pdf with structured content instead."
        ) from e

    _t = _HTML_TITLE_RE.search(html)
    _h = _HTML_H1_RE.search(html)
    filename = _derive_filename(
        filename, ".pdf",
        title=_html_text(_t.group(1)) if _t else None,
        hints=(
            _html_text(_h.group(1)) if _h else None,
            _first_words(_html_text(html)),
        ),
    )
    buf = io.BytesIO()
    # SSRF/LFI guard (audit-2026 re-audit round 9): WeasyPrint's DEFAULT
    # url_fetcher resolves file:// and http(s) URLs embedded in the (possibly
    # injection-controlled) HTML. `<img src="file:///proc/self/environ">` would
    # embed the agent's env secrets into the PDF; `<img
    # src="http://169.254.169.254/…">` or an internal host would SSRF. Restrict
    # to inline data: URIs + public https only.
    HTML(string=html or "", url_fetcher=_safe_pdf_url_fetcher).write_pdf(buf)
    return await _persist(buf.getvalue(), filename, MIME_PDF, user_scope)


def _safe_pdf_url_fetcher(url: str):
    """Allow only inline data: URIs and public https for WeasyPrint PDF
    rendering (audit-2026 re-audit round 9). Blocks file:// (LFI) and any
    private/loopback/link-local/metadata host (SSRF)."""
    from weasyprint.urls import default_url_fetcher  # type: ignore
    u = (url or "").strip()
    low = u.lower()
    if low.startswith("data:"):
        return default_url_fetcher(u)
    if low.startswith("https://"):
        from app.agent.smart_fetch.reader import _assert_public_url
        _assert_public_url(u)  # raises ValueError on private/internal/metadata
        return default_url_fetcher(u)
    raise ValueError(f"blocked non-public URL in PDF generation: {u[:80]!r}")
